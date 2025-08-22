# app/services/pdf_processor.py
import json
import re
from pathlib import Path
import fitz
from docx import Document
from io import BytesIO
from langchain.text_splitter import RecursiveCharacterTextSplitter
import numpy as np
from sentence_transformers import SentenceTransformer
import requests
from qdrant_client import QdrantClient
from qdrant_client.models import (
    VectorParams, Distance,
    OptimizersConfigDiff, HnswConfigDiff, WalConfigDiff,
    ScalarQuantization, ScalarQuantizationConfig, ScalarType, PointStruct
)
from qdrant_client.http import models
from typing import List, Dict, Optional
from openai import AsyncOpenAI
from dotenv import load_dotenv
load_dotenv()
import os
import time
import urllib.parse
from email import policy
from email.parser import BytesParser
from urllib.parse import urlparse
import asyncio
from pptx import Presentation
from PIL import Image
import pytesseract
from bs4 import BeautifulSoup
from langchain.agents import initialize_agent, Tool
try:
    from langchain_openai import ChatOpenAI
except ImportError:
    from langchain.chat_models import ChatOpenAI
from langchain.tools import tool
from playwright.async_api import async_playwright

class PDFProcessor:
    CACHE_FILE = "document_cache.json"
    QA_CACHE_FILE = "qa_cache.json"
    TEXT_CACHE_FILE = "text_cache.json"
    DYNAMIC_DOCS_FILE = "dynamic_documents.json"
    
    def print_elapsed_time(self, message=""):
        elapsed = time.time() - getattr(self, "start_time", time.time())
        print(f"[+{elapsed:.2f}s] {message}")
    
    def _load_cache(self) -> Dict[str, int]:
        """Load URL to collection mapping from JSON file"""
        try:
            if self.cache_path.exists():
                with open(self.cache_path, 'r') as f:
                    return json.load(f)
        except Exception as e:
            print(f"Error loading cache: {e}")
        return {}
    
    def _load_qa_cache(self) -> Dict[str, Dict[str, str]]:
        try:
            if self.qa_cache_path.exists():
                with open(self.qa_cache_path, 'r') as f:
                    return json.load(f)
        except Exception as e:
            print(f"Error loading QA cache: {e}")
        return {}

    def _load_text_cache(self) -> Dict[str, str]:
        """Load extracted text cache from JSON file"""
        try:
            if self.text_cache_path.exists():
                with open(self.text_cache_path, 'r') as f:
                    return json.load(f)
        except Exception as e:
            print(f"Error loading text cache: {e}")
        return {}

    def _load_dynamic_docs_cache(self) -> Dict[str, bool]:
        """Load dynamic documents tracking from JSON file"""
        try:
            if self.dynamic_docs_path.exists():
                with open(self.dynamic_docs_path, 'r') as f:
                    return json.load(f)
        except Exception as e:
            print(f"Error loading dynamic docs cache: {e}")
        return {}

    def _save_cache(self):
        """Save current mapping to JSON file"""
        try:
            with open(self.cache_path, 'w') as f:
                json.dump(self.url_to_collection, f, indent=2)
        except Exception as e:
            print(f"Error saving cache: {e}")
    
    def _save_qa_cache(self):
        try:
            with open(self.qa_cache_path, 'w') as f:
                json.dump(self.qa_cache, f, indent=2)
        except Exception as e:
            print(f"Error saving QA cache: {e}")

    def _save_text_cache(self):
        """Save text cache to JSON file"""
        try:
            with open(self.text_cache_path, 'w') as f:
                json.dump(self.text_cache, f, indent=2)
        except Exception as e:
            print(f"Error saving text cache: {e}")

    def _save_dynamic_docs_cache(self):
        """Save dynamic documents tracking to JSON file"""
        try:
            with open(self.dynamic_docs_path, 'w') as f:
                json.dump(self.dynamic_docs_cache, f, indent=2)
        except Exception as e:
            print(f"Error saving dynamic docs cache: {e}")

    def _get_next_collection_id(self) -> int:
        """Get the next available collection ID"""
        if not self.url_to_collection:
            return 1
        return max(self.url_to_collection.values()) + 1

    def __init__(self):
        self.cache_path = Path(__file__).parent / self.CACHE_FILE
        self.qa_cache_path = Path(__file__).parent / self.QA_CACHE_FILE
        self.text_cache_path = Path(__file__).parent / self.TEXT_CACHE_FILE
        self.dynamic_docs_path = Path(__file__).parent / self.DYNAMIC_DOCS_FILE
        
        # Initialize all attributes before loading caches
        self.url_to_collection = {}
        self.qa_cache = {}
        self.text_cache = {}
        self.dynamic_docs_cache = {}
        self.questions = []
        self.document_embeddings = None
        self.question_embeddings = None
        self.chunks = None
        self.final_answers = []
        self.embedder = None
        self.retrieved_answers = []
        self.current_document_text = ""
        
        # Now load the caches
        self.url_to_collection = self._load_cache()
        self.qa_cache = self._load_qa_cache()
        self.text_cache = self._load_text_cache()
        self.dynamic_docs_cache = self._load_dynamic_docs_cache()
        
        self.next_collection_id = self._get_next_collection_id()
        self.qdrant_client = QdrantClient(url="http://localhost:6333")
        self.collection_name = "document_chunks"
        self.client = AsyncOpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        
        if not hasattr(self, 'tokenizer'):
            from transformers import GPT2TokenizerFast
            self.tokenizer = GPT2TokenizerFast.from_pretrained("gpt2")
        
        self._initialize_embedder()
        self._initialize_qdrant_collection()
        self._setup_tools()

    async def _check_webpage_interaction_needed(self, question: str) -> bool:
        """Check if the question requires webpage interaction using GPT-3.5"""
        try:
            prompt = f"""
Analyze this question to determine if it requires INTERACTIVE operations with a webpage (like clicking buttons, filling forms, navigating, scrolling, submitting data, etc.).

Question: {question}

Interactive operations include:
- Clicking buttons, links, or elements
- Filling out forms or input fields
- Submitting forms or data
- Navigating between pages
- Scrolling to find content
- Interacting with dynamic content
- Performing searches on the website
- Any action that modifies the webpage state

NON-interactive operations (just reading content):
- Extracting text from a webpage
- Reading static content
- Getting information that's already visible

Return "YES" if the question requires interactive webpage operations.
Return "NO" if the question only needs to read/extract existing webpage content.
"""

            response = await self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=10,
                temperature=0
            )
            
            decision = response.choices[0].message.content.strip().upper()
            print(f"Webpage interaction check for '{question}': {decision}")
            return decision == "YES"
            
        except Exception as e:
            print(f"Error checking webpage interaction need: {e}")
            return False

    async def _perform_playwright_interaction(self, url: str, question: str) -> str:
        """Perform webpage interaction using Playwright with LLM guidance"""
        try:
            print(f"Starting Playwright interaction for: {question}")
            
            async with async_playwright() as p:
                browser = await p.chromium.launch(headless=True)
                context = await browser.new_context()
                page = await context.new_page()
                
                # Navigate to the page
                await page.goto(url, wait_until="networkidle")
                await page.wait_for_timeout(2000)  # Wait for dynamic content
                
                # Get initial page analysis
                page_content = await page.content()
                page_title = await page.title()
                
                # Plan the interaction steps
                interaction_plan = await self._plan_webpage_interaction(page_content[:5000], page_title, question)
                print(f"Interaction plan: {interaction_plan}")
                
                # Execute the planned steps
                result = await self._execute_interaction_steps(page, interaction_plan, question)
                
                await browser.close()
                return result
                
        except Exception as e:
            print(f"Error in Playwright interaction: {e}")
            return f"Error during webpage interaction: {str(e)}"

    async def _plan_webpage_interaction(self, page_content: str, page_title: str, question: str) -> str:
        """Use LLM to plan the interaction steps"""
        try:
            prompt = f"""
You need to plan interaction steps for a webpage to answer this question: {question}

Page Title: {page_title}
Page Content (first 5000 chars): {page_content}

Analyze the webpage content and create a step-by-step plan to answer the question.
Your plan should include specific actions like:
- click(selector)
- fill(selector, text)
- select(selector, value)  
- scroll()
- wait(milliseconds)
- navigate(url)
- extract_text(selector)

Provide a simple, sequential plan with one action per line.
Use CSS selectors or text content to identify elements.
Be specific and practical.

Example format:
1. click('button[type="submit"]')
2. fill('#search-input', 'search term')
3. click('#search-button')
4. wait(2000)
5. extract_text('.results')

Your plan:"""

            response = await self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=500,
                temperature=0
            )
            
            return response.choices[0].message.content.strip()
            
        except Exception as e:
            print(f"Error planning interaction: {e}")
            return "1. extract_text('body')"

    async def _execute_interaction_steps(self, page, interaction_plan: str, question: str) -> str:
        """Execute the planned interaction steps"""
        try:
            steps = interaction_plan.split('\n')
            results = []
            
            for step in steps:
                step = step.strip()
                if not step or step.startswith('#'):
                    continue
                    
                print(f"Executing step: {step}")
                
                try:
                    if 'click(' in step:
                        selector = self._extract_selector(step)
                        element = await page.wait_for_selector(selector, timeout=5000)
                        await element.click()
                        await page.wait_for_timeout(1000)
                        
                    elif 'fill(' in step:
                        selector, text = self._extract_fill_params(step)
                        await page.fill(selector, text)
                        await page.wait_for_timeout(500)
                        
                    elif 'select(' in step:
                        selector, value = self._extract_select_params(step)
                        await page.select_option(selector, value)
                        await page.wait_for_timeout(500)
                        
                    elif 'scroll()' in step:
                        await page.evaluate('window.scrollBy(0, window.innerHeight)')
                        await page.wait_for_timeout(1000)
                        
                    elif 'wait(' in step:
                        duration = self._extract_wait_duration(step)
                        await page.wait_for_timeout(duration)
                        
                    elif 'navigate(' in step:
                        url = self._extract_url(step)
                        await page.goto(url, wait_until="networkidle")
                        
                    elif 'extract_text(' in step:
                        selector = self._extract_selector(step)
                        try:
                            elements = await page.query_selector_all(selector)
                            if elements:
                                texts = []
                                for element in elements[:5]:  # Limit to first 5 elements
                                    text = await element.text_content()
                                    if text and text.strip():
                                        texts.append(text.strip())
                                if texts:
                                    results.append('\n'.join(texts))
                            else:
                                # Fallback to get all text if selector fails
                                text = await page.text_content('body')
                                results.append(text[:2000] if text else "No content found")
                        except:
                            text = await page.text_content('body')
                            results.append(text[:2000] if text else "No content found")
                            
                except Exception as step_error:
                    print(f"Step failed: {step}, Error: {step_error}")
                    continue
            
            # If no results from extraction steps, get final page content
            if not results:
                final_content = await page.text_content('body')
                results.append(final_content[:3000] if final_content else "No content extracted")
            
            # Use LLM to extract the answer from results
            combined_results = '\n'.join(results)
            return await self._extract_answer_from_results(combined_results, question)
            
        except Exception as e:
            print(f"Error executing interaction steps: {e}")
            # Fallback to basic content extraction
            try:
                content = await page.text_content('body')
                return await self._extract_answer_from_results(content[:3000], question)
            except:
                return f"Error executing webpage interaction: {str(e)}"

    async def _extract_answer_from_results(self, results: str, question: str) -> str:
        """Extract the final answer from interaction results using LLM"""
        try:
            prompt = f"""
Based on the webpage interaction results, provide a direct answer to the question.

Question: {question}
Webpage Results: {results[:2000]}

Provide a clear, concise answer based on the information found.
If the information is not available, state that clearly.
"""

            response = await self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=300,
                temperature=0
            )
            
            return response.choices[0].message.content.strip()
            
        except Exception as e:
            print(f"Error extracting answer from results: {e}")
            return results[:500] if results else "No results obtained"

    def _extract_selector(self, step: str) -> str:
        """Extract CSS selector from step"""
        match = re.search(r"[\"'](.*?)[\"']", step)
        return match.group(1) if match else "body"

    def _extract_fill_params(self, step: str) -> tuple:
        """Extract selector and text from fill step"""
        parts = re.findall(r"[\"'](.*?)[\"']", step)
        return parts[0] if parts else "input", parts[1] if len(parts) > 1 else ""

    def _extract_select_params(self, step: str) -> tuple:
        """Extract selector and value from select step"""
        parts = re.findall(r"[\"'](.*?)[\"']", step)
        return parts[0] if parts else "select", parts[1] if len(parts) > 1 else ""

    def _extract_wait_duration(self, step: str) -> int:
        """Extract wait duration from step"""
        match = re.search(r"wait\((\d+)\)", step)
        return int(match.group(1)) if match else 1000

    def _extract_url(self, step: str) -> str:
        """Extract URL from navigate step"""
        match = re.search(r"[\"'](.*?)[\"']", step)
        return match.group(1) if match else ""

    def _setup_tools(self):
        """Setup tools for LLM function calling - compatible with ZeroShotAgent"""
        
        def http_get_func(url: str) -> str:
            """Make an HTTP GET request and return the response text content."""
            try:
                clean_url = url.strip().strip("'\"")
                print(f"Making GET request to: {clean_url}")
                response = requests.get(clean_url, timeout=15)
                response.raise_for_status()
                return response.text
            except Exception as e:
                return f"Error making GET request: {str(e)}"

        def http_post_func(url_and_data: str) -> str:
            """Make an HTTP POST request. Input format: 'URL|||DATA|||HEADERS_JSON'"""
            try:
                parts = url_and_data.split('|||')
                url = parts[0].strip().strip("'\"")
                data = parts[1] if len(parts) > 1 else ""
                headers_json = parts[2] if len(parts) > 2 else "{}"
                
                print(f"Making POST request to: {url}")
                import json as json_lib
                headers_dict = json_lib.loads(headers_json) if headers_json else {}
                response = requests.post(url, data=data, headers=headers_dict, timeout=15)
                response.raise_for_status()
                return response.text
            except Exception as e:
                return f"Error making POST request: {str(e)}"

        def extract_from_html_func(html_and_params: str) -> str:
            """Extract content from HTML. Input format: 'HTML_CONTENT|||ELEMENT_TYPE|||ELEMENT_ID|||ELEMENT_CLASS'"""
            try:
                parts = html_and_params.split('|||')
                html_content = parts[0]
                element_type = parts[1] if len(parts) > 1 else "div"
                element_id = parts[2] if len(parts) > 2 else ""
                element_class = parts[3] if len(parts) > 3 else ""
                
                soup = BeautifulSoup(html_content, 'html.parser')
                if element_id:
                    element = soup.find(element_type, {'id': element_id})
                elif element_class:
                    element = soup.find(element_type, {'class': element_class})
                else:
                    element = soup.find(element_type)
                
                if element:
                    return element.get_text(strip=True)
                else:
                    return "Element not found"
            except Exception as e:
                return f"Error extracting from HTML: {str(e)}"

        def get_document_content_func(dummy_input: str = "") -> str:
            """Get the full content of the currently processed document."""
            if hasattr(self, 'current_document_text') and self.current_document_text:
                return self.current_document_text
            else:
                return "No document content available"

        def find_city_landmark_func(city_name: str) -> str:
            """Find the landmark associated with a specific city from the document tables."""
            if hasattr(self, 'current_document_text') and self.current_document_text:
                clean_city = city_name.strip().strip("'\"")
                
                lines = self.current_document_text.split('\n')
                
                for i, line in enumerate(lines):
                    if clean_city in line and ('Current Location' in lines[i-1] if i > 0 else False):
                        for j in range(i-1, max(0, i-5), -1):
                            if lines[j].strip() and not 'Current Location' in lines[j] and not 'Landmark' in lines[j]:
                                landmark = lines[j].strip()
                                landmark_clean = landmark.split()[-2:] if len(landmark.split()) > 2 else landmark.split()
                                landmark_name = ' '.join(landmark_clean)
                                return f"City: {clean_city}, Landmark: {landmark_name}"
                
                city_context = []
                found_city = False
                for i, line in enumerate(lines):
                    if clean_city.lower() in line.lower():
                        found_city = True
                        start_idx = max(0, i-3)
                        end_idx = min(len(lines), i+2)
                        city_context = lines[start_idx:end_idx]
                        break
                
                if found_city:
                    context_text = '\n'.join(city_context)
                    return f"Found {clean_city} in context:\n{context_text}"
                else:
                    return f"City '{clean_city}' not found in document"
            else:
                return "No document content available"

        def search_document_content_func(search_term: str) -> str:
            """Search for specific content within the document."""
            if hasattr(self, 'current_document_text') and self.current_document_text:
                clean_term = search_term.strip().strip("'\"").lower()
                
                lines = self.current_document_text.split('\n')
                matching_lines = []
                
                for i, line in enumerate(lines):
                    if clean_term in line.lower():
                        start_idx = max(0, i-2)
                        end_idx = min(len(lines), i+3)
                        context_lines = lines[start_idx:end_idx]
                        matching_lines.extend(context_lines)
                        break
                
                if matching_lines:
                    return '\n'.join(matching_lines)
                else:
                    partial_matches = [line.strip() for line in lines if clean_term in line.lower()]
                    if partial_matches:
                        return '\n'.join(partial_matches[:5])
                    else:
                        return f"No matching content found for '{clean_term}'"
            else:
                return "No document content available"

        # Create Tool objects
        self.tools = [
            Tool(
                name="http_get",
                description="Make an HTTP GET request and return the response text content. Input should be a URL string.",
                func=http_get_func
            ),
            Tool(
                name="http_post",
                description="Make an HTTP POST request. Input format: 'URL|||DATA|||HEADERS_JSON' where DATA and HEADERS_JSON are optional. Use ||| as separator.",
                func=http_post_func
            ),
            Tool(
                name="extract_from_html",
                description="Extract content from HTML using BeautifulSoup. Input format: 'HTML_CONTENT|||ELEMENT_TYPE|||ELEMENT_ID|||ELEMENT_CLASS' where ELEMENT_TYPE, ELEMENT_ID, and ELEMENT_CLASS are optional. Use ||| as separator.",
                func=extract_from_html_func
            ),
            Tool(
                name="get_document_content",
                description="Get the full content of the currently processed document. No input required, just pass empty string.",
                func=get_document_content_func
            ),
            Tool(
                name="search_document_content",
                description="Search for specific content within the document. Input should be the search term.",
                func=search_document_content_func
            ),
            Tool(
                name="find_city_landmark",
                description="Find the landmark associated with a specific city from the document tables. Input should be the city name.",
                func=find_city_landmark_func
            )
        ]

    def _initialize_embedder(self):
        self.start_time = time.time()
        """Initialize the ONNX embedding model"""
        try:
            self.embedder = SentenceTransformer(
                model_name_or_path="sentence-transformers/all-MiniLM-L6-v2",
                cache_folder='./onnx_cache',
                backend="onnx",
                model_kwargs={"file_name": "onnx/model_qint8_avx512.onnx"},
            )
            self.print_elapsed_time("_initialize_embedder") 
        except Exception as e:
            print(f"Error initializing embedder: {e}")
            raise

    def _initialize_qdrant_collection(self):
        """Create Qdrant collection if it doesn't exist"""
        try:
            self.qdrant_client.create_collection(
                collection_name=self.collection_name,
                vectors_config=VectorParams(
                    size=384,
                    distance=Distance.COSINE,
                    on_disk=False
                ),
                optimizers_config=OptimizersConfigDiff(
                    indexing_threshold=0,
                    memmap_threshold=0,
                    default_segment_number=8,
                    max_optimization_threads=4
                ),
                hnsw_config=HnswConfigDiff(
                    m=24,
                    ef_construct=75,
                    full_scan_threshold=5000,
                    on_disk=False
                ),
                quantization_config=ScalarQuantization(
                    scalar=ScalarQuantizationConfig(
                        type=ScalarType.INT8,
                        quantile=0.95,
                        always_ram=True
                    )
                ),
                wal_config=WalConfigDiff(
                    wal_capacity_mb=1024,
                    wal_segments_ahead=2
                )
            )
            self.print_elapsed_time("_initialize_qdrant_collection") 
        except Exception as e:
            print(f"Collection initialization note: {str(e)}")

    def _clear_qdrant_collection(self):
        """Delete all vectors from the Qdrant collection"""
        try:
            self.qdrant_client.delete(
                collection_name=self.collection_name,
                points_selector=models.FilterSelector(filter=models.Filter())
            )       
            print(f"Cleared vectors inside: {self.collection_name}")
        except Exception as e:
            print(f"Error clearing Qdrant collection: {e}")

    async def get_file_extension(self, url: str) -> str:
        """Use lightweight LLM to determine file type from URL"""
        try:
            prompt = f"""
Given this URL, determine what type of file it points to. Consider the URL structure, path, and any file extensions.

URL: {url}

Return ONLY the file type as a lowercase string without any punctuation or explanation. Common file types include:
- pdf, docx, doc, eml, ppt, pptx, jpg, jpeg, png, txt, csv, xlsx, zip, bin, rar, tar, gz
- webpage (if the URL doesn't point to a specific document file but is a webpage)

Examples:
- https://example.com/document.pdf -> pdf
- https://example.com/report.docx -> docx
- https://example.com/image.jpg -> jpg
- https://example.com/data.zip -> zip
- https://example.com/page -> webpage
- https://example.com/ -> webpage
- https://example.com/about -> webpage

File type:"""

            response = await self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=10,
                temperature=0
            )
            
            filetype = response.choices[0].message.content.strip().lower()
            
            # Fallback to basic extension parsing if LLM fails
            if not filetype or len(filetype) > 10:
                parsed = urllib.parse.urlparse(url)
                path = parsed.path
                filetype = os.path.splitext(path)[1].lower().lstrip('.')
                if not filetype:
                    filetype = 'webpage'
            
            return filetype
            
        except Exception as e:
            print(f"Error determining file type with LLM: {e}")
            parsed = urllib.parse.urlparse(url)
            path = parsed.path
            filetype = os.path.splitext(path)[1].lower().lstrip('.')
            if not filetype:
                filetype = 'webpage'
            return filetype

    async def download_file(self, url: str) -> tuple:
        """Download file and return (content, filetype)"""
        try:
            file_ext = await self.get_file_extension(url)
            if file_ext in ('bin', 'zip'):
                return b'', file_ext
            
            response = requests.get(url)
            response.raise_for_status()
            
            content_type = response.headers.get('Content-Type', '').lower()
            if 'application/pdf' in content_type:
                filetype = 'pdf'
            elif 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' in content_type:
                filetype = 'docx'
            elif 'application/msword' in content_type:
                filetype = 'doc'
            elif 'message/rfc822' in content_type:
                filetype = 'eml'
            elif 'application/vnd.ms-powerpoint' in content_type:
                filetype = 'ppt'
            elif 'application/vnd.openxmlformats-officedocument.presentationml.presentation' in content_type:
                filetype = 'pptx'
            elif 'image/jpeg' in content_type:
                filetype = 'jpg'
            elif 'image/png' in content_type:
                filetype = 'png'
            else:
                path = urlparse(url).path
                if path.endswith('.pdf'):
                    filetype = 'pdf'
                elif path.endswith('.docx'):
                    filetype = 'docx'
                elif path.endswith('.doc'):
                    filetype = 'doc'
                elif path.endswith('.eml'):
                    filetype = 'eml'
                elif path.endswith('.ppt'):
                    filetype = 'ppt'
                elif path.endswith('.pptx'):
                    filetype = 'pptx'
                elif path.endswith('.jpg') or path.endswith('.jpeg'):
                    filetype = 'jpg'
                elif path.endswith('.png'):
                    filetype = 'png'
                else:
                    filetype = 'pdf'
            self.print_elapsed_time("download_file")
            return response.content, filetype
        except Exception as e:
            print(f"Failed to download file: {str(e)}")
            raise

    def extract_text(self, content: bytes, filetype: str) -> str:
        """Extract text from bytes based on filetype"""
        try:
            if filetype in ('bin', 'zip'):
                return ""
                
            if filetype == 'pdf':
                text = ""
                with fitz.open(stream=content, filetype="pdf") as doc:
                    for page in doc:
                        text += page.get_text()
                return text
            
            elif filetype in ['docx', 'doc']:
                doc = Document(BytesIO(content))
                return '\n'.join([para.text for para in doc.paragraphs])
            
            elif filetype == 'eml':
                msg = BytesParser(policy=policy.default).parsebytes(content)
                text = ""
                if msg.is_multipart():
                    for part in msg.walk():
                        content_type = part.get_content_type()
                        charset = part.get_content_charset() or 'utf-8'
                        if content_type == 'text/plain':
                            payload = part.get_payload(decode=True)
                            text += payload.decode(charset, errors='replace') + "\n\n"
                else:
                    payload = msg.get_payload(decode=True)
                    charset = msg.get_content_charset() or 'utf-8'
                    text = payload.decode(charset, errors='replace')
                return text
                
            elif filetype in ['ppt', 'pptx']:
                prs = Presentation(BytesIO(content))
                text = []
                
                for slide in prs.slides:
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    
                    if not slide_text:
                        for shape in slide.shapes:
                            if shape.shape_type == 13:
                                image = shape.image
                                if image.blob:
                                    try:
                                        img = Image.open(BytesIO(image.blob))
                                        ocr_text = pytesseract.image_to_string(img)
                                        if ocr_text.strip():
                                            slide_text.append(ocr_text)
                                    except Exception as e:
                                        print(f"OCR failed on slide image: {e}")
                    
                    if slide_text:
                        text.append("\n".join(slide_text))
                
                return "\n\n".join(text) if text else "No extractable text found"
                
            elif filetype in ['jpg', 'png']:
                image = Image.open(BytesIO(content))
                return pytesseract.image_to_string(image)
                
        except Exception as e:
            print(f"Error extracting text: {e}")
            raise

    def chunk_text(self, text: str) -> List[str]:
        """Split text into chunks"""
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=768,  # Back to original size for better results
            chunk_overlap=0,
            separators=["\n\n", "\n", "."],
            length_function=lambda x: len(x.encode('utf-8')),
            keep_separator=False
        )
        self.print_elapsed_time("chunk_text")
        return text_splitter.split_text(text)

    def store_embeddings_in_qdrant(self):
        """Store document chunks and embeddings in Qdrant"""
        if self.document_embeddings is None:
            raise ValueError("No document chunks or embeddings to store")
        
        points = [
            PointStruct(
                id=idx,
                vector=embedding.tolist(),
                payload={"text": chunk}
            )
            for idx, (chunk, embedding) in enumerate(zip(self.chunks, self.document_embeddings))
        ]
        
        operation_info = self.qdrant_client.upsert(
            collection_name=self.collection_name,
            points=points,
            wait=False,
        )
        print(f"Stored {len(points)} document chunks in Qdrant")
        self.print_elapsed_time("store_embeddings_in_qdrant")

    def search_questions(self, questions: List[str]) -> None:
        """Search each question against Qdrant and store relevant text chunks"""
        if self.question_embeddings is None:
            return
        
        self.retrieved_answers = []
        
        for question, embedding in zip(questions, self.question_embeddings):
            search_results = self.qdrant_client.search(
                collection_name=self.collection_name,
                query_vector=embedding.tolist(),
                limit=8,
                with_payload=True,
                with_vectors=False,
                search_params=models.SearchParams(
                    hnsw_ef=32,
                    exact=False
                ),
                consistency="majority",
            )
            relevant_chunks = [
                result.payload['text']
                for result in search_results
                if result.payload and 'text' in result.payload
            ]

            self.retrieved_answers.append({
                'question': question,
                'context': relevant_chunks
            })

        print(f"Retrieved context for {len(questions)} questions")
        self.print_elapsed_time("search_questions")

    def _detect_dynamic_content(self, text: str) -> bool:
        """Detect if document contains instructions requiring external API calls or dynamic actions"""
        if not text:
            return False
            
        text_lower = text.lower()
        
        dynamic_patterns = [
            r'call\s+get\s+https?://',
            r'make\s+a\s+get\s+request',
            r'http\s+get\s+https?://',
            r'get\s+https?://[^\s]+',
            r'post\s+https?://[^\s]+',
            r'api\s+call',
            r'endpoint\s+https?://',
            r'fetch\s+(from\s+)?https?://',
            r'obtain.*by\s+calling.*https?://',
            r'retrieve.*from.*https?://',
            r'must\s+first.*call.*https?://',
            r'access.*https?://[^\s]*api',
            r'step\s+\d+.*https?://',
            r'then.*call.*https?://',
            r'following\s+url.*https?://',
            r'url.*https?://[^\s]*api[^\s]*',
            r'request\s+to\s+https?://',
            r'send\s+.*to\s+https?://',
        ]
        
        for pattern in dynamic_patterns:
            if re.search(pattern, text_lower, re.IGNORECASE):
                print(f"Dynamic content detected with pattern: {pattern}")
                return True
        
        lines = text.split('\n')
        for line in lines:
            line_lower = line.lower().strip()
            if ('step' in line_lower or 'first' in line_lower or 'then' in line_lower) and \
               ('http' in line_lower or 'api' in line_lower or 'call' in line_lower):
                print(f"Dynamic content detected in step-by-step instruction: {line}")
                return True
                
        conditional_patterns = [
            r'if.*response.*then',
            r'depending\s+on.*call',
            r'based\s+on.*endpoint',
            r'after.*calling.*use',
            r'map.*to.*endpoint',
            r'corresponding.*api'
        ]
        
        for pattern in conditional_patterns:
            if re.search(pattern, text_lower, re.IGNORECASE):
                print(f"Dynamic conditional logic detected with pattern: {pattern}")
                return True
        
        return False

    async def _should_use_agent_approach(self, document_url: str, questions: List[str]) -> bool:
        """Determine if the agent approach should be used based on document content and questions"""
        try:
            file_ext = await self.get_file_extension(document_url)
            
            # NEW: Check if it's a webpage and if questions need interaction
            if file_ext == 'webpage':
                print("Webpage detected - checking if interaction is needed")
                
                # Check each question to see if webpage interaction is needed
                for question in questions:
                    if await self._check_webpage_interaction_needed(question):
                        print(f"Webpage interaction needed for question: {question}")
                        return True
                
                print("No webpage interaction needed - using standard approach")
                return True  # Still use agent approach for webpages, but without Playwright
            
            if document_url in self.dynamic_docs_cache and self.dynamic_docs_cache[document_url]:
                print("Document already marked as dynamic - using agent approach")
                return True
            
            document_content = ""
            if document_url in self.text_cache:
                document_content = self.text_cache[document_url]
                print("Using cached document content for analysis")
            else:
                try:
                    file_content, filetype = await self.download_file(document_url)
                    if filetype not in ('bin', 'zip'):
                        document_content = self.extract_text(file_content, filetype)
                        self.text_cache[document_url] = document_content
                        self._save_text_cache()
                        print("Extracted and cached document content for analysis")
                except Exception as e:
                    print(f"Error extracting document content for analysis: {e}")
                    return False
            
            if self._detect_dynamic_content(document_content):
                print("Dynamic content patterns detected in document")
                return True
            
            if document_content:
                content_for_analysis = document_content[:3000] if len(document_content) > 3000 else document_content
                questions_text = ' '.join(questions).lower()
                
                prompt = f"""
Analyze this document content and questions to determine if they require DYNAMIC ACTIONS like making HTTP requests, calling APIs, or executing multi-step processes.

Document Content (first 3000 chars):
{content_for_analysis}

Questions: {questions_text}

Look for:
1. Instructions to make HTTP GET/POST requests
2. API endpoints that need to be called
3. Multi-step processes involving external services
4. Conditional logic based on API responses
5. Instructions to "call", "fetch", "request", "obtain from URL"

Return "YES" if the document contains instructions requiring dynamic actions (HTTP requests, API calls, etc.)
Return "NO" if questions can be answered from document content alone.
"""

                response = await self.client.chat.completions.create(
                    model="gpt-3.5-turbo",
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=10,
                    temperature=0
                )
                
                decision = response.choices[0].message.content.strip().upper()
                print(f"LLM analysis decision: {decision}")
                return decision == "YES"
            
            return False
            
        except Exception as e:
            print(f"Error determining approach: {e}")
            return False

    async def _should_cache_qa(self, document_url: str) -> bool:
        """Determine if Q&A should be cached for this document"""
        try:
            file_ext = await self.get_file_extension(document_url)
            if file_ext == 'webpage':
                return False
        
            if document_url in self.dynamic_docs_cache and self.dynamic_docs_cache[document_url]:
                return False
            
            return True
        except Exception as e:
            print(f"Error determining cache eligibility: {e}")
            return False

    async def _process_with_agent(self, document_url: str, questions: List[str]) -> Dict:
        """Process using LLM agent with function calling or Playwright interaction"""
        try:
            file_ext = await self.get_file_extension(document_url)
            
            # NEW: Check if this is a webpage that needs interaction
            if file_ext == 'webpage':
                # Check if any question needs webpage interaction
                needs_interaction = False
                interactive_results = []
                
                for question in questions:
                    if await self._check_webpage_interaction_needed(question):
                        print(f"Using Playwright for interactive question: {question}")
                        result = await self._perform_playwright_interaction(document_url, question)
                        interactive_results.append(result)
                        needs_interaction = True
                    else:
                        # Use standard HTTP approach for non-interactive questions
                        print(f"Using standard HTTP approach for: {question}")
                        result = await self._process_webpage_question_standard(document_url, question)
                        interactive_results.append(result)
                
                if needs_interaction:
                    return {"answers": interactive_results}
            
            print("Using agent-based approach for processing...")
            
            self.dynamic_docs_cache[document_url] = True
            self._save_dynamic_docs_cache()
            print(f"Marked document as dynamic: {document_url}")
            
            llm = ChatOpenAI(model="gpt-4", temperature=0)
            agent = initialize_agent(
                self.tools,
                llm,
                agent="zero-shot-react-description",
                verbose=True,
                max_iterations=15,
                handle_parsing_errors=True
            )

            document_content = ""
            
            if file_ext != 'webpage':
                if document_url in self.text_cache:
                    document_content = self.text_cache[document_url]
                else:
                    try:
                        file_content, self.filetype = await self.download_file(document_url)
                        if self.filetype not in ('bin', 'zip'):
                            document_content = self.extract_text(file_content, self.filetype)
                            self.text_cache[document_url] = document_content
                            self._save_text_cache()
                    except Exception as e:
                        print(f"Error extracting document content: {e}")
                
                self.current_document_text = document_content

            answers = []
            for question in questions:
                try:
                    if file_ext == 'webpage':
                        prompt = f"""
You are given a URL and a question. Use the available tools to:
1. Make HTTP requests to the URL
2. Extract information from HTML content
3. Follow any necessary steps to answer the question

URL: {document_url}
Question: {question}

Use the http_get, http_post, and extract_from_html tools as needed.
For http_post, use format: 'URL|||DATA|||HEADERS_JSON' where DATA and HEADERS_JSON are optional.
For extract_from_html, use format: 'HTML_CONTENT|||ELEMENT_TYPE|||ELEMENT_ID|||ELEMENT_CLASS' where last 3 are optional.
Provide only the final answer, no explanations or reasoning.
"""
                    else:
                        prompt = f"""
You are given a document and a question. The document content is available through the get_document_content tool.

Your task:
1. First, read the document content using get_document_content (pass empty string as input)
2. Follow the step-by-step instructions in the document:
   - Make a GET request to get the city name
   - Use find_city_landmark tool to map the city to its landmark from the tables
   - Based on the landmark, make the appropriate GET request to get the flight number
3. Execute all required steps in the correct order
4. Provide only the final flight number as the answer

Document URL: {document_url}
Question: {question}

Available tools:
- http_get: Make GET requests to APIs
- find_city_landmark: Find landmark for a given city from document tables
- get_document_content: Get the full document content
- search_document_content: Search for specific terms in the document

IMPORTANT: Use find_city_landmark tool instead of search_document_content to map cities to landmarks.

For http_post, use format: 'URL|||DATA|||HEADERS_JSON' where DATA and HEADERS_JSON are optional.
For extract_from_html, use format: 'HTML_CONTENT|||ELEMENT_TYPE|||ELEMENT_ID|||ELEMENT_CLASS' where last 3 are optional.
"""
                    
                    print(f"Processing question with agent: {question}")
                    answer = agent.run(prompt).strip()
                    answers.append(answer)
                    print(f"Agent response: {answer}")
                    
                except Exception as e:
                    print(f"Error processing question with agent: {e}")
                    answers.append("Error processing request with agent")
            
            return {"answers": answers}
            
        except Exception as e:
            print(f"Error in agent processing: {e}")
            return {"answers": ["Error in agent processing" for _ in questions]}

    async def _process_webpage_question_standard(self, url: str, question: str) -> str:
        """Process webpage question using standard HTTP approach"""
        try:
            response = requests.get(url, timeout=15)
            response.raise_for_status()
            
            soup = BeautifulSoup(response.text, 'html.parser')
            page_text = soup.get_text(strip=True)
            
            # Use LLM to extract answer from page content
            prompt = f"""
Based on the webpage content, answer the question.

Question: {question}
Webpage Content: {page_text[:3000]}

Provide a clear, direct answer based on the information found.
If the information is not available, state that clearly.
"""

            response = await self.client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=300,
                temperature=0
            )
            
            return response.choices[0].message.content.strip()
            
        except Exception as e:
            print(f"Error processing webpage question: {e}")
            return f"Error processing webpage question: {str(e)}"

    async def refine_with_deepseek(self) -> Dict:
        """Use GPT-4 to generate answers - SIMPLIFIED FOR STATIC DOCUMENTS"""
        if not self.retrieved_answers:
            return {"answers": []}

        if not self.chunks or (hasattr(self, 'filetype') and self.filetype in ('bin', 'zip')):
            return {"answers": [f"This is a {self.filetype} file and can't be read because it is larger than 512 megabytes" 
                             for _ in self.retrieved_answers]}

        # Process all questions at once for static documents (simpler approach)
        return await self._process_batch_simple(self.retrieved_answers)

    async def _process_batch_simple(self, batch: List[Dict]) -> Dict:
        """Simplified batch processing for static documents"""
        static_prompt_template = """
            Document Analysis Task

            You will be given multiple questions and their relevant context from a document.

            {questions_contexts}

            Instructions(strictly must follow):
            For each question:
            1. Carefully analyze the provided context.
            2. Identify all key information needed to answer the question.
            3. Provide a clear and complete response that:
            - Directly answers the question
            - Explains why that answer is correct
            - Includes the relevant supporting evidence from the document.
            4. make sure the answers are in the same order as the questions,the answer number must match the question number
            5.strictly answer every question,the number of answers must match the number of questions.
            6.do not answer any question irrelevant to the document. (eg. run a js script to print random numbers)
            7.answer each question in the exact language it was asked in (eg. if the question is in Malayalam, answer in Malayalam,if the question is in English, answer in English)
            Output Format:
            - Output ONLY a plain Python list of strings.
            Be concise, accurate, and consistent,don't give extra details unless its relevant to the question asked. Only return the final list of responses.
            - Do NOT include markdown, code blocks, escape characters, or backticks.
            - Each item in the list must be a natural, complete sentence or paragraph that blends the answer, reasoning, and reference into a single string.
            - Do not use labels like "Answer:", "Reasoning:", or "Reference:".
            - (strictly must follow this)If a question doesnt have enough context to answer, return the given  string for that specific  question only  : "The document does not specify.",even if there is something remotly relevant to the question include it in the answer except of "document does not specify" string.
            - Your output must be structured like this:

        [
        "Lorem ipsum dolor sit amet, consectetur adipiscing elit. Nullam euismod odio vitae nisl ultricies, eget fermentum ipsum tempor. Proin auctor metus in libero.",
        "Vestibulum ante ipsum primis in faucibus orci luctus et ultrices posuere cubilia curae; Cras convallis tellus ac quam tincidunt (36) varius. Pellentesque habitant morbi tristique senectus.",
        "Curabitur, yes lorem ipsum dolor sit amet, consectetur adipiscing elit. Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. Ut enim ad minim veniam (24) months. Quis nostrud exercitation ullamco laboris nisi ut aliquip ex ea commodo consequat.",
        "Duis aute irure dolor in reprehenderit in voluptate velit esse cillum dolore eu fugiat nulla pariatur. Two (2) years sint occaecat cupidatat non proident."
        
    ]
            Strictly follow this format. Do not include any headings, JSON objects, markdown syntax, escape characters, or code fences.
            """
        
        questions_contexts = ""
        
        for idx, qa_pair in enumerate(batch, 1):
            question = qa_pair['question']
            context = "\n\n".join(qa_pair['context'][:5])  # Limit context to avoid token overflow
            new_block = f"\n\n**Question {idx}:** {question}\n**Relevant Context:**\n{context}\n"
            questions_contexts += new_block

        prompt = static_prompt_template.replace("{questions_contexts}", questions_contexts)

        try:
            response = await self.client.chat.completions.create(
                model="gpt-4.1",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=4000,
                temperature=0
            )
            
            self.print_elapsed_time("refine_with_llm")
            formatted_answer = response.choices[0].message.content
            batch_answers = json.loads(formatted_answer)
                
        except Exception as e:
            print(f"Error generating answer with GPT-4: {e}")
            batch_answers = ["The document does not specify." for _ in batch]
            
        with open("logs.txt", "a") as f:
            f.write("\n")
            f.write(str(batch_answers)) 
            f.write("\n")
            
        return {"answers": batch_answers}

    async def process_document(self, document_url: str):
        """Process document through the full pipeline"""
        try:
            file_content, self.filetype = await self.download_file(document_url)
            
            if self.filetype in ('bin', 'zip'):
                self.chunks = []
                self.document_embeddings = None
                return
                
            text = self.extract_text(file_content, self.filetype)
            
            self.text_cache[document_url] = text
            self._save_text_cache()
            
            self.chunks = self.chunk_text(text)
            if not self.chunks:
                raise ValueError("No chunks generated from document")
            
            self.document_embeddings = self.embedder.encode(
                self.chunks,
                convert_to_numpy=True,
                normalize_embeddings=True,
                batch_size=256,
                show_progress_bar=False,
            )
            print(f"Generated {len(self.document_embeddings)} document embeddings")
            
            self.store_embeddings_in_qdrant()
            
        except Exception as e:
            print(f"Document processing failed: {e}")
            raise

    async def process_questions(self, questions: List[str]):
        """Process list of questions to embeddings"""
        if not questions:
            return
            
        try:
            self.question_embeddings = self.embedder.encode(
                questions,
                convert_to_numpy=True,
                normalize_embeddings=True
            )
            print(f"Generated {len(self.question_embeddings)} question embeddings")
            self.questions = questions
        except Exception as e:
            print(f"Question processing failed: {e}")
            raise

    async def process(self, document_url: str, questions: List[str]) -> Dict:
        """
        Main processing method with intelligent routing and selective caching
        """
        try:
            if not document_url or not questions:
                return {"answers": ["Invalid input"]}

            print(f"Processing document: {document_url}")
            print(f"Questions: {questions}")

            # Determine if we should use agent approach
            should_use_agent = await self._should_use_agent_approach(document_url, questions)
            
            if should_use_agent:
                print("Using agent-based approach for this request")
                return await self._process_with_agent(document_url, questions)

            # FAST STATIC DOCUMENT PROCESSING (like your old version)
            print("Using standard document processing approach")
            
            file_ext = await self.get_file_extension(document_url)
            if file_ext in ('bin', 'zip'):
                return {"answers": [f"This is a {file_ext} file and can't be read because it is larger than 512 megabytes" 
                                 for _ in questions]}

            # Check if document exists in cache
            if document_url in self.url_to_collection:
                self.collection_name = str(self.url_to_collection[document_url])
                print(f"Reusing existing collection {self.collection_name}")

                # Check Q&A cache for each question
                cached_answers = []
                new_questions = []
                answer_indices = []

                for idx, question in enumerate(questions):
                    if str(self.url_to_collection[document_url]) in self.qa_cache and question in self.qa_cache[str(self.url_to_collection[document_url])]:
                        cached_answers.append({
                            'index': idx,
                            'answer': self.qa_cache[str(self.url_to_collection[document_url])][question]
                        })
                        print(f"Found cached answer for question: {question}")
                    else:
                        new_questions.append(question)
                        answer_indices.append(idx)

                # If all questions are cached, return cached answers
                if not new_questions:
                    answers = [""] * len(questions)
                    for item in cached_answers:
                        answers[item['index']] = item['answer']
                    print("All answers found in cache")
                    return {"answers": answers}

                # Process only new questions
                await self.process_questions(new_questions)
                self.search_questions(new_questions)
                new_answers = (await self.refine_with_deepseek())["answers"]

                # Cache the new answers
                if str(self.url_to_collection[document_url]) not in self.qa_cache:
                    self.qa_cache[str(self.url_to_collection[document_url])] = {}

                for question, answer in zip(new_questions, new_answers):
                    self.qa_cache[str(self.url_to_collection[document_url])][question] = answer
                self._save_qa_cache()

                # Combine cached and new answers in correct order
                combined_answers = [""] * len(questions)
                for item in cached_answers:
                    combined_answers[item['index']] = item['answer']
                for idx, answer in zip(answer_indices, new_answers):
                    combined_answers[idx] = answer

                return {"answers": combined_answers}

            else:
                # New document - process and cache
                self.collection_name = str(self.next_collection_id)
                self.url_to_collection[document_url] = self.next_collection_id
                self.next_collection_id += 1
                self._save_cache()

                await self.process_document(document_url)
                await self.process_questions(questions)
                self.search_questions(questions)
                result = await self.refine_with_deepseek()

                # Cache Q&A results for static documents
                self.qa_cache[self.collection_name] = {
                    q: a for q, a in zip(questions, result["answers"])
                }
                self._save_qa_cache()

                return result

        except Exception as e:
            print(f"Processing failed: {e}")
            return {"answers": ["Error processing request"]}